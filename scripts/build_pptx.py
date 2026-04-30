#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py — JSON spec → .pptx 演示文稿

依赖: python-pptx (pip install python-pptx)

使用:
    python build_pptx.py --spec deck.json --out my_deck.pptx
    python build_pptx.py --list-themes
    python build_pptx.py --list-layouts
    python build_pptx.py --schema  # 打印 JSON spec 模板

JSON Spec:
{
  "title": "标题（可选，会变成第一张 title 幻灯片）",
  "subtitle": "副标题",
  "theme": "modern|minimal|tech|warm|blueprint",   # 默认 modern
  "footer": "底部小字（可选）",
  "slides": [
    {"layout": "title",       "title": "...", "subtitle": "..."},
    {"layout": "section",     "title": "第 X 章"},
    {"layout": "content",     "title": "...", "bullets": ["...", "..."]},
    {"layout": "two_column",  "title": "...", "left_title": "...", "left_bullets": [...], "right_title": "...", "right_bullets": [...]},
    {"layout": "quote",       "text": "...", "author": "..."},
    {"layout": "image",       "title": "...", "image": "path/to/img.png", "caption": "..."},
    {"layout": "table",       "title": "...", "headers": ["A","B","C"], "rows": [["..","..",".."], ...]},
    {"layout": "stats",       "title": "...", "stats": [{"value": "85%", "label": "良品率"}, ...]},
    {"layout": "thanks",      "title": "Thank You", "subtitle": "..."}
  ]
}

设计原则:
  - 大字号、强对比、足量留白
  - 单页只放一个核心信息
  - 12-16 张为佳；超过 25 张自动 warning
"""

import sys
import os
import io
import json
import argparse
from pathlib import Path

# 强制 UTF-8 stdout（防 Windows GBK 乱码）
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
else:
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError as e:
    print(f"[build_pptx] ERROR: python-pptx not installed.\n  pip install python-pptx\n  Detail: {e}", file=sys.stderr)
    sys.exit(2)

# ---- Themes -----------------------------------------------------------------
# 五套手调过的颜色方案。每个都给：bg / fg / accent / muted / surface
THEMES = {
    'modern': {
        'bg':       (0xF7, 0xF8, 0xFA),
        'fg':       (0x1A, 0x1F, 0x2C),
        'accent':   (0x5B, 0x8D, 0xEF),  # blue
        'accent2':  (0xF6, 0xA0, 0x46),  # warm orange for highlights
        'muted':    (0x6B, 0x73, 0x80),
        'surface':  (0xFF, 0xFF, 0xFF),
        'title_font': 'Microsoft YaHei',
        'body_font':  'Microsoft YaHei',
    },
    'minimal': {
        'bg':       (0xFF, 0xFF, 0xFF),
        'fg':       (0x11, 0x11, 0x11),
        'accent':   (0x00, 0x00, 0x00),
        'accent2':  (0xE5, 0x3E, 0x3E),
        'muted':    (0x99, 0x99, 0x99),
        'surface':  (0xF5, 0xF5, 0xF5),
        'title_font': 'Microsoft YaHei',
        'body_font':  'Microsoft YaHei',
    },
    'tech': {
        'bg':       (0x0B, 0x0E, 0x14),
        'fg':       (0xE6, 0xE9, 0xEF),
        'accent':   (0x4F, 0xC3, 0xF7),
        'accent2':  (0x80, 0xCB, 0xC4),
        'muted':    (0x8B, 0x96, 0xA8),
        'surface':  (0x16, 0x1B, 0x24),
        'title_font': 'Microsoft YaHei',
        'body_font':  'Microsoft YaHei',
    },
    'warm': {
        'bg':       (0xFD, 0xF6, 0xEC),
        'fg':       (0x3A, 0x2A, 0x1F),
        'accent':   (0xC2, 0x6E, 0x4A),
        'accent2':  (0x6B, 0x8E, 0x4E),
        'muted':    (0x8B, 0x77, 0x65),
        'surface':  (0xFF, 0xFF, 0xFF),
        'title_font': 'Microsoft YaHei',
        'body_font':  'Microsoft YaHei',
    },
    'blueprint': {
        'bg':       (0x1F, 0x3A, 0x68),
        'fg':       (0xF0, 0xF6, 0xFF),
        'accent':   (0xFF, 0xC8, 0x4D),
        'accent2':  (0x9B, 0xC1, 0xFF),
        'muted':    (0x9C, 0xB3, 0xD4),
        'surface':  (0x29, 0x4A, 0x80),
        'title_font': 'Microsoft YaHei',
        'body_font':  'Microsoft YaHei',
    },
}
LAYOUTS = ['title', 'section', 'content', 'two_column', 'quote', 'image', 'table', 'stats', 'thanks']

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(t):
    return RGBColor(*t)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_textbox(slide, left, top, w, h, text, *, font_size=18, bold=False,
                color=(0, 0, 0), font='Microsoft YaHei', align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = font
    return tb


def add_accent_bar(slide, theme, top=Inches(0.7), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(0.6), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(theme['accent'])
    bar.line.fill.background()


def add_footer(slide, theme, footer_text, page_no=None):
    if footer_text:
        add_textbox(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
                    footer_text, font_size=10, color=theme['muted'],
                    font=theme['body_font'])
    if page_no is not None:
        add_textbox(slide, Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.35),
                    str(page_no), font_size=10, color=theme['muted'],
                    font=theme['body_font'], align=PP_ALIGN.RIGHT)


def add_paragraph_runs(tf, items, *, font_size=20, color=(0, 0, 0), font='Microsoft YaHei',
                       bullet=False, line_spacing=1.3, accent_first=None):
    """Replace text_frame content with given items (list of strings)."""
    if not items:
        return
    tf.clear()
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if bullet:
            run = p.add_run()
            run.text = '•  '
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(accent_first if accent_first else color)
            run.font.name = font
            run2 = p.add_run()
            run2.text = str(t)
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = rgb(color)
            run2.font.name = font
        else:
            run = p.add_run()
            run.text = str(t)
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(color)
            run.font.name = font


# ---- Layout builders --------------------------------------------------------

def slide_title(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, theme['bg'])
    # large accent block on left
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(2.5),
                                   Inches(0.4), Inches(2.5))
    block.fill.solid()
    block.fill.fore_color.rgb = rgb(theme['accent'])
    block.line.fill.background()
    add_textbox(slide, Inches(1.0), Inches(2.7), Inches(11), Inches(1.6),
                slide_data.get('title', ''), font_size=54, bold=True,
                color=theme['fg'], font=theme['title_font'])
    if slide_data.get('subtitle'):
        add_textbox(slide, Inches(1.0), Inches(4.4), Inches(11), Inches(0.8),
                    slide_data['subtitle'], font_size=22,
                    color=theme['muted'], font=theme['body_font'])
    add_footer(slide, theme, footer, page_no)


def slide_section(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['accent'])
    add_textbox(slide, Inches(0.6), Inches(3.0), Inches(12), Inches(1.5),
                slide_data.get('title', ''), font_size=60, bold=True,
                color=(0xFF, 0xFF, 0xFF), font=theme['title_font'],
                align=PP_ALIGN.CENTER)
    if slide_data.get('subtitle'):
        add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
                    slide_data['subtitle'], font_size=20,
                    color=(0xFF, 0xFF, 0xFF), font=theme['body_font'],
                    align=PP_ALIGN.CENTER)


def slide_content(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['bg'])
    add_textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.6),
                slide_data.get('title', ''), font_size=30, bold=True,
                color=theme['fg'], font=theme['title_font'])
    add_accent_bar(slide, theme, top=Inches(1.2))
    bullets = slide_data.get('bullets', [])
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    add_paragraph_runs(tf, bullets, font_size=22, color=theme['fg'],
                       font=theme['body_font'], bullet=True,
                       accent_first=theme['accent'], line_spacing=1.45)
    add_footer(slide, theme, footer, page_no)


def slide_two_column(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['bg'])
    add_textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.6),
                slide_data.get('title', ''), font_size=30, bold=True,
                color=theme['fg'], font=theme['title_font'])
    add_accent_bar(slide, theme, top=Inches(1.2))
    for i, side in enumerate(['left', 'right']):
        x = Inches(0.6 + i * 6.2)
        title = slide_data.get(f'{side}_title', '')
        if title:
            add_textbox(slide, x, Inches(1.5), Inches(5.8), Inches(0.5),
                        title, font_size=20, bold=True,
                        color=theme['accent'], font=theme['title_font'])
        body = slide.shapes.add_textbox(x, Inches(2.1), Inches(5.8), Inches(4.7))
        tf = body.text_frame
        tf.word_wrap = True
        bullets = slide_data.get(f'{side}_bullets', [])
        add_paragraph_runs(tf, bullets, font_size=18, color=theme['fg'],
                           font=theme['body_font'], bullet=True,
                           accent_first=theme['accent'])
    add_footer(slide, theme, footer, page_no)


def slide_quote(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['bg'])
    # giant quote mark
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(2),
                '"', font_size=120, bold=True,
                color=theme['accent'], font='Georgia')
    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(11), Inches(2.8),
                slide_data.get('text', ''), font_size=32, bold=False,
                color=theme['fg'], font=theme['title_font'])
    if slide_data.get('author'):
        add_textbox(slide, Inches(1.5), Inches(5.4), Inches(11), Inches(0.6),
                    f"— {slide_data['author']}", font_size=18,
                    color=theme['muted'], font=theme['body_font'])
    add_footer(slide, theme, footer, page_no)


def slide_image(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['bg'])
    if slide_data.get('title'):
        add_textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.6),
                    slide_data['title'], font_size=30, bold=True,
                    color=theme['fg'], font=theme['title_font'])
        add_accent_bar(slide, theme, top=Inches(1.2))
        img_top, img_h = Inches(1.5), Inches(5.0)
    else:
        img_top, img_h = Inches(0.7), Inches(5.8)
    img_path = slide_data.get('image')
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2.5), img_top,
                                 width=Inches(8.3), height=img_h)
    else:
        # placeholder rectangle
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), img_top,
                                    Inches(8.3), img_h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = rgb(theme['surface'])
        ph.line.color.rgb = rgb(theme['muted'])
        add_textbox(slide, Inches(2.5), img_top + Inches(2.5), Inches(8.3), Inches(0.6),
                    f"[图片缺失: {img_path or '未提供'}]", font_size=14,
                    color=theme['muted'], font=theme['body_font'],
                    align=PP_ALIGN.CENTER)
    if slide_data.get('caption'):
        add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
                    slide_data['caption'], font_size=12,
                    color=theme['muted'], font=theme['body_font'],
                    align=PP_ALIGN.CENTER)
    add_footer(slide, theme, footer, page_no)


def slide_table(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['bg'])
    add_textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.6),
                slide_data.get('title', ''), font_size=30, bold=True,
                color=theme['fg'], font=theme['title_font'])
    add_accent_bar(slide, theme, top=Inches(1.2))

    headers = slide_data.get('headers', [])
    rows = slide_data.get('rows', [])
    if not headers and not rows:
        return
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    n_rows = (1 if headers else 0) + len(rows)
    if n_cols == 0 or n_rows == 0:
        return
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.6), Inches(1.6),
        Inches(12.1), Inches(min(5.3, 0.5 + n_rows * 0.5))
    )
    table = table_shape.table
    if headers:
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(theme['accent'])
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.color.rgb = rgb((0xFF, 0xFF, 0xFF))
                    r.font.size = Pt(14)
                    r.font.name = theme['body_font']
    row_offset = 1 if headers else 0
    for i, row in enumerate(rows):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i + row_offset, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = rgb(theme['fg'])
                    r.font.name = theme['body_font']
    add_footer(slide, theme, footer, page_no)


def slide_stats(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['bg'])
    add_textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.6),
                slide_data.get('title', ''), font_size=30, bold=True,
                color=theme['fg'], font=theme['title_font'])
    add_accent_bar(slide, theme, top=Inches(1.2))

    stats = slide_data.get('stats', [])[:4]
    if not stats:
        return
    n = len(stats)
    total_w = 12.1
    gap = 0.3
    card_w = (total_w - gap * (n - 1)) / n
    for i, s in enumerate(stats):
        x = Inches(0.6 + i * (card_w + gap))
        # card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, Inches(2.0), Inches(card_w), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(theme['surface'])
        card.line.color.rgb = rgb(theme['muted'])
        # value
        add_textbox(slide, x, Inches(2.4), Inches(card_w), Inches(1.6),
                    str(s.get('value', '')), font_size=56, bold=True,
                    color=theme['accent'], font=theme['title_font'],
                    align=PP_ALIGN.CENTER)
        # label
        add_textbox(slide, x, Inches(4.3), Inches(card_w), Inches(1.0),
                    str(s.get('label', '')), font_size=15,
                    color=theme['fg'], font=theme['body_font'],
                    align=PP_ALIGN.CENTER)
    add_footer(slide, theme, footer, page_no)


def slide_thanks(prs, slide_data, theme, footer, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, theme['accent'])
    add_textbox(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(2),
                slide_data.get('title', 'Thank You'), font_size=78, bold=True,
                color=(0xFF, 0xFF, 0xFF), font=theme['title_font'],
                align=PP_ALIGN.CENTER)
    if slide_data.get('subtitle'):
        add_textbox(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.8),
                    slide_data['subtitle'], font_size=22,
                    color=(0xFF, 0xFF, 0xFF), font=theme['body_font'],
                    align=PP_ALIGN.CENTER)


LAYOUT_BUILDERS = {
    'title':      slide_title,
    'section':    slide_section,
    'content':    slide_content,
    'two_column': slide_two_column,
    'quote':      slide_quote,
    'image':      slide_image,
    'table':      slide_table,
    'stats':      slide_stats,
    'thanks':     slide_thanks,
}


def build(spec, out_path):
    theme_name = spec.get('theme', 'modern')
    if theme_name not in THEMES:
        print(f"[build_pptx] unknown theme '{theme_name}', falling back to 'modern'", file=sys.stderr)
        theme_name = 'modern'
    theme = THEMES[theme_name]
    footer = spec.get('footer', '')

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = list(spec.get('slides', []))
    # auto-prepend title slide if top-level title is set and first slide isn't already a title
    if spec.get('title') and not (slides and slides[0].get('layout') == 'title'):
        slides.insert(0, {
            'layout': 'title',
            'title': spec['title'],
            'subtitle': spec.get('subtitle', ''),
        })

    if not slides:
        print("[build_pptx] ERROR: spec has no slides", file=sys.stderr)
        sys.exit(2)

    if len(slides) > 25:
        print(f"[build_pptx] WARN: {len(slides)} slides — consider trimming for impact (sweet spot is 12-16)", file=sys.stderr)

    for i, sd in enumerate(slides):
        layout = sd.get('layout', 'content')
        if layout not in LAYOUT_BUILDERS:
            print(f"[build_pptx] WARN slide {i+1}: unknown layout '{layout}', defaulting to 'content'", file=sys.stderr)
            layout = 'content'
        LAYOUT_BUILDERS[layout](prs, sd, theme, footer, i + 1)

    prs.save(out_path)


def cli():
    ap = argparse.ArgumentParser(description="JSON spec → .pptx 演示文稿")
    ap.add_argument('--spec', help="JSON spec 文件路径")
    ap.add_argument('--stdin', action='store_true', help="从 stdin 读 JSON spec")
    ap.add_argument('--out', default='deck.pptx', help="输出 .pptx 路径")
    ap.add_argument('--list-themes', action='store_true')
    ap.add_argument('--list-layouts', action='store_true')
    ap.add_argument('--schema', action='store_true', help="打印 JSON spec 模板")
    args = ap.parse_args()

    if args.list_themes:
        print('Themes:', ', '.join(THEMES.keys()))
        return
    if args.list_layouts:
        print('Layouts:', ', '.join(LAYOUTS))
        return
    if args.schema:
        sample = {
            "title": "示例 PPT",
            "subtitle": "你的副标题",
            "theme": "modern",
            "footer": "Manuscopy · 2026",
            "slides": [
                {"layout": "section", "title": "第一部分"},
                {"layout": "content", "title": "要点", "bullets": ["要点 1", "要点 2", "要点 3"]},
                {"layout": "two_column", "title": "对比",
                 "left_title": "方案 A", "left_bullets": ["优点 1", "优点 2"],
                 "right_title": "方案 B", "right_bullets": ["优点 1", "优点 2"]},
                {"layout": "stats", "title": "数据",
                 "stats": [{"value": "85%", "label": "良品率"},
                           {"value": "1.2s", "label": "节拍"},
                           {"value": "3x", "label": "效率"}]},
                {"layout": "quote", "text": "工艺即细节，细节即标准。", "author": "某老工艺员"},
                {"layout": "table", "title": "参数表",
                 "headers": ["材料", "进给", "转速"],
                 "rows": [["6061-T6", "0.15", "8000"], ["45钢", "0.10", "1200"]]},
                {"layout": "thanks", "title": "感谢观看", "subtitle": "Q&A"}
            ]
        }
        print(json.dumps(sample, ensure_ascii=False, indent=2))
        return

    if args.stdin:
        spec = json.load(sys.stdin)
    elif args.spec:
        with open(args.spec, 'r', encoding='utf-8') as f:
            spec = json.load(f)
    else:
        ap.error("provide --spec FILE or --stdin (or use --schema to see the format)")

    out_path = Path(args.out).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    build(spec, str(out_path))
    size = out_path.stat().st_size
    print(f"OK: {out_path} ({size/1024:.1f} KB, {len(spec.get('slides', []))} slides)")


if __name__ == '__main__':
    cli()
